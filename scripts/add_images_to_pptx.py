#!/usr/bin/env python3
"""
Agrega imágenes (diagramas, QR codes) a la presentación PPTX
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

# Rutas
BASE_DIR = Path(__file__).parent.parent
PPTX_PATH = BASE_DIR / 'output' / 'slides-nwc-labitconf.pptx'
OUTPUT_PATH = BASE_DIR / 'output' / 'slides-nwc-labitconf.pptx'

# Imágenes a insertar
IMAGES = {
    4: {  # Slide 4: "Qué es NWC" - diagrama de flujo simple
        'path': BASE_DIR / 'assets' / 'nwc-simple-flow.png',
        'left': Inches(7),
        'top': Inches(3.5),
        'width': Inches(5.5),
    },
    11: {  # Slide 11: "IA + MCP" - arquitectura NWC+MCP
        'path': BASE_DIR / 'assets' / 'architecture-nwc-mcp.png',
        'left': Inches(1),
        'top': Inches(4),
        'width': Inches(8.5),
    },
    12: {  # Slide 12: "Cómo Empezar Hoy" - QR codes
        'qr_codes': [
            {
                'path': BASE_DIR / 'assets' / 'qr-nip47.png',
                'left': Inches(1.5),
                'top': Inches(4.5),
                'width': Inches(1.5),
            },
            {
                'path': BASE_DIR / 'assets' / 'qr-alby.png',
                'left': Inches(4),
                'top': Inches(4.5),
                'width': Inches(1.5),
            },
            {
                'path': BASE_DIR / 'assets' / 'qr-nostr.png',
                'left': Inches(6.5),
                'top': Inches(4.5),
                'width': Inches(1.5),
            },
        ]
    }
}


def add_images_to_presentation():
    """Agrega imágenes a la presentación"""
    if not PPTX_PATH.exists():
        print(f"❌ Error: No se encontró {PPTX_PATH}")
        return

    print("🖼️  Agregando imágenes a la presentación...")
    prs = Presentation(str(PPTX_PATH))

    for slide_num, image_data in IMAGES.items():
        slide_idx = slide_num - 1  # 0-indexed

        if slide_idx >= len(prs.slides):
            print(f"⚠️  Slide {slide_num} no existe, saltando...")
            continue

        slide = prs.slides[slide_idx]

        # Caso especial: múltiples QR codes
        if 'qr_codes' in image_data:
            for qr in image_data['qr_codes']:
                if qr['path'].exists():
                    slide.shapes.add_picture(
                        str(qr['path']),
                        qr['left'],
                        qr['top'],
                        width=qr['width']
                    )
                    print(f"  ✅ Agregado QR a Slide {slide_num}: {qr['path'].name}")
                else:
                    print(f"  ⚠️  No encontrado: {qr['path']}")
        else:
            # Imagen única
            if image_data['path'].exists():
                slide.shapes.add_picture(
                    str(image_data['path']),
                    image_data['left'],
                    image_data['top'],
                    width=image_data['width']
                )
                print(f"  ✅ Agregado a Slide {slide_num}: {image_data['path'].name}")
            else:
                print(f"  ⚠️  No encontrado: {image_data['path']}")

    # Guardar
    prs.save(str(OUTPUT_PATH))

    # Verificar tamaño
    file_size_mb = OUTPUT_PATH.stat().st_size / (1024 * 1024)
    print(f"\n✅ Presentación actualizada: {OUTPUT_PATH}")
    print(f"📊 Tamaño: {file_size_mb:.2f} MB {'✅' if file_size_mb < 15 else '⚠️  EXCEDE 15MB'}")

    if file_size_mb >= 15:
        print("⚠️  El archivo excede el límite de 15MB. Considera optimizar imágenes.")


if __name__ == '__main__':
    add_images_to_presentation()
