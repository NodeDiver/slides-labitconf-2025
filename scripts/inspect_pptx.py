#!/usr/bin/env python3
"""
Inspects the generated PPTX to see what content is actually there
"""

from pathlib import Path
from pptx import Presentation

OUTPUT_PATH = Path(__file__).parent.parent / 'output' / 'slides-nwc-labitconf.pptx'

def inspect_pptx():
    """Inspect the PPTX file"""
    if not OUTPUT_PATH.exists():
        print(f"❌ File not found: {OUTPUT_PATH}")
        return

    prs = Presentation(str(OUTPUT_PATH))

    print(f"📊 Total slides: {len(prs.slides)}")
    print()

    for idx, slide in enumerate(prs.slides, 1):
        print(f"{'='*60}")
        print(f"SLIDE {idx}")
        print(f"{'='*60}")
        print(f"Layout: {slide.slide_layout.name}")

        # Check title
        if slide.shapes.title:
            print(f"Title: {slide.shapes.title.text if slide.shapes.title.text else '(empty)'}")
        else:
            print("Title: (no title placeholder)")

        # Check all shapes
        print(f"\nShapes ({len(slide.shapes)}):")
        for shape_idx, shape in enumerate(slide.shapes):
            print(f"  [{shape_idx}] Type: {shape.shape_type}, ", end="")
            if hasattr(shape, 'name'):
                print(f"Name: {shape.name}, ", end="")
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    print(f"Text: {text[:80]}...")
                else:
                    print(f"Text: (empty)")
            else:
                print("(no text)")

        # Check notes
        notes_text = slide.notes_slide.notes_text_frame.text
        if notes_text.strip():
            print(f"\nNotes: {notes_text[:100]}...")

        print()

if __name__ == '__main__':
    inspect_pptx()
