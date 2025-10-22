#!/usr/bin/env python3
"""
Inspects the template to see available layouts and their placeholders
"""

from pathlib import Path
from pptx import Presentation

TEMPLATE_PATH = Path(__file__).parent.parent / 'Template_SPEAKERS_2025.pptx'

def inspect_template():
    """Inspect the template file"""
    if not TEMPLATE_PATH.exists():
        print(f"❌ File not found: {TEMPLATE_PATH}")
        return

    prs = Presentation(str(TEMPLATE_PATH))

    print(f"📊 Total slides in template: {len(prs.slides)}")
    print(f"📐 Total layouts: {len(prs.slide_layouts)}")
    print()

    for idx, layout in enumerate(prs.slide_layouts):
        print(f"{'='*60}")
        print(f"LAYOUT {idx}: {layout.name}")
        print(f"{'='*60}")
        print(f"Placeholders ({len(layout.placeholders)}):")

        for placeholder in layout.placeholders:
            print(f"  - idx={placeholder.placeholder_format.idx}, type={placeholder.placeholder_format.type}, name={placeholder.name}")

        print(f"\nShapes ({len(layout.shapes)}):")
        for shape in layout.shapes:
            if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                print(f"  [PLACEHOLDER] {shape.name} (type={shape.placeholder_format.type})")
            else:
                print(f"  [SHAPE] {shape.name} (type={shape.shape_type})")

        print()

if __name__ == '__main__':
    inspect_template()
