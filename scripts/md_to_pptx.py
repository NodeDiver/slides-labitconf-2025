#!/usr/bin/env python3
"""
Convierte el contenido Markdown de slides-nwc.md a PPTX usando el template de LABITCONF.
"""

import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Rutas
BASE_DIR = Path(__file__).parent.parent
TEMPLATE_PATH = BASE_DIR / 'Template_SPEAKERS_2025.pptx'
CONTENT_PATH = BASE_DIR / 'content' / 'slides-nwc.md'
OUTPUT_PATH = BASE_DIR / 'output' / 'slides-nwc-labitconf.pptx'

# Colores del template LABITCONF
COLORS = {
    'accent1': RGBColor(0x42, 0x85, 0xF4),  # Azul
    'accent2': RGBColor(0x21, 0x21, 0x21),  # Negro
    'accent3': RGBColor(0x78, 0x90, 0x9C),  # Gris azulado
    'accent4': RGBColor(0xFF, 0xAB, 0x40),  # Naranja
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'black': RGBColor(0x00, 0x00, 0x00),
}

# Mapeo de layouts (basado en análisis del template)
LAYOUT_MAP = {
    'cover': 0,      # Cover - Color 01
    'content': 4,    # TITLE_AND_BODY 2
    'two-column': 5, # TITLE_AND_TWO_COLUMNS
    'section': 3,    # SECTION_HEADER
}


class SlideParser:
    """Parser para el formato Markdown de slides"""

    def __init__(self, md_content):
        self.content = md_content
        self.slides = []

    def parse(self):
        """Parsea el Markdown y extrae información de cada slide"""
        # Dividir por separadores ---
        sections = re.split(r'\n---\n', self.content)

        for section in sections:
            if not section.strip() or section.startswith('#'):
                continue  # Saltar encabezado y secciones vacías

            slide_data = self._parse_slide_section(section)
            if slide_data:
                self.slides.append(slide_data)

        return self.slides

    def _parse_slide_section(self, section):
        """Parsea una sección individual de slide"""
        lines = section.strip().split('\n')
        slide_data = {
            'layout': 4,  # Default: TITLE_AND_BODY 2
            'title': '',
            'bullets': [],
            'notes': '',
            'metadata': {}
        }

        in_notes = False
        in_bullets = False
        current_section = None

        for line in lines:
            line_stripped = line.strip()

            # Extraer metadata de comentarios
            if line_stripped.startswith('**Slide'):
                continue
            elif line_stripped.startswith('- Layout:'):
                layout_match = re.search(r'Layout:\s*(\d+)', line_stripped)
                if layout_match:
                    slide_data['layout'] = int(layout_match.group(1)) - 1  # 0-indexed
            elif line_stripped.startswith('- Type:'):
                type_match = re.search(r'Type:\s*(\w+)', line_stripped)
                if type_match:
                    layout_type = type_match.group(1)
                    slide_data['layout'] = LAYOUT_MAP.get(layout_type, 4)

            # Título (##)
            elif line_stripped.startswith('##'):
                slide_data['title'] = line_stripped.lstrip('#').strip()

            # Subtítulos (###) para dos columnas
            elif line_stripped.startswith('###'):
                current_section = line_stripped.lstrip('#').strip()
                if 'columns' not in slide_data:
                    slide_data['columns'] = {}
                slide_data['columns'][current_section] = []

            # Bullets
            elif line_stripped.startswith('-') or line_stripped.startswith('•'):
                bullet_text = line_stripped.lstrip('-•').strip()
                if bullet_text.startswith('**'):
                    bullet_text = bullet_text  # Mantener formato bold

                if current_section and 'columns' in slide_data:
                    slide_data['columns'][current_section].append(bullet_text)
                else:
                    slide_data['bullets'].append(bullet_text)

            # Notas
            elif line_stripped.startswith('**Notas'):
                in_notes = True
                notes_text = line_stripped.replace('**Notas**:', '').replace('**Notas**', '').strip()
                slide_data['notes'] = notes_text
            elif in_notes and line_stripped:
                slide_data['notes'] += ' ' + line_stripped

            # Pendientes (marcadores para completar)
            elif '[Pendiente:' in line_stripped:
                slide_data['metadata']['pending'] = line_stripped

        return slide_data if slide_data['title'] else None


def create_presentation(slides_data):
    """Crea la presentación PPTX usando el template"""
    prs = Presentation(str(TEMPLATE_PATH))

    # Limpiar slides de ejemplo del template (conservar solo layouts)
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    # Crear slides a partir de los datos parseados
    for slide_data in slides_data:
        layout_idx = slide_data.get('layout', 4)
        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)

        # Agregar título
        if slide.shapes.title:
            slide.shapes.title.text = slide_data['title']

        # Agregar contenido según el tipo de layout
        if 'columns' in slide_data:
            _add_two_column_content(slide, slide_data['columns'])
        elif slide_data['bullets']:
            _add_bullets(slide, slide_data['bullets'])

        # Agregar notas del orador
        if slide_data['notes']:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data['notes']

    return prs


def _add_bullets(slide, bullets):
    """Agrega bullets al slide"""
    # Buscar el placeholder de body
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            text_frame = shape.text_frame
            text_frame.clear()  # Limpiar contenido de ejemplo

            for idx, bullet in enumerate(bullets):
                p = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
                p.text = _clean_markdown(bullet)
                p.level = 0
                p.font.size = Pt(24)  # Tamaño legible
            break


def _add_two_column_content(slide, columns_data):
    """Agrega contenido en dos columnas"""
    # Buscar placeholders de body (debería haber 2 para layout de dos columnas)
    body_shapes = [s for s in slide.shapes if s.has_text_frame and s != slide.shapes.title]

    column_idx = 0
    for section_title, bullets in columns_data.items():
        if column_idx >= len(body_shapes):
            break

        text_frame = body_shapes[column_idx].text_frame
        text_frame.clear()

        # Agregar subtítulo de la columna
        p = text_frame.paragraphs[0]
        p.text = section_title
        p.font.bold = True
        p.font.size = Pt(28)

        # Agregar bullets
        for bullet in bullets:
            p = text_frame.add_paragraph()
            p.text = _clean_markdown(bullet)
            p.level = 0
            p.font.size = Pt(22)

        column_idx += 1


def _clean_markdown(text):
    """Limpia formato Markdown simple (**, *, [], etc.)"""
    # Remover bold
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    # Remover cursiva
    text = re.sub(r'\*(.*?)\*', r'\1', text)
    # Remover links [text](url)
    text = re.sub(r'\[(.*?)\]\(.*?\)', r'\1', text)
    return text


def main():
    """Función principal"""
    print("🚀 Generando presentación PPTX desde Markdown...")

    # Verificar que existan los archivos necesarios
    if not TEMPLATE_PATH.exists():
        print(f"❌ Error: No se encontró el template en {TEMPLATE_PATH}")
        return

    if not CONTENT_PATH.exists():
        print(f"❌ Error: No se encontró el contenido en {CONTENT_PATH}")
        return

    # Leer contenido Markdown
    with open(CONTENT_PATH, 'r', encoding='utf-8') as f:
        md_content = f.read()

    # Parsear slides
    parser = SlideParser(md_content)
    slides_data = parser.parse()

    print(f"📄 Parseadas {len(slides_data)} slides del Markdown")

    # Crear presentación
    prs = create_presentation(slides_data)

    # Guardar
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))

    # Verificar tamaño
    file_size_mb = OUTPUT_PATH.stat().st_size / (1024 * 1024)
    print(f"✅ Presentación generada: {OUTPUT_PATH}")
    print(f"📊 Tamaño: {file_size_mb:.2f} MB {'✅' if file_size_mb < 15 else '⚠️  EXCEDE 15MB'}")

    if file_size_mb >= 15:
        print("⚠️  El archivo excede el límite de 15MB. Considera optimizar imágenes.")


if __name__ == '__main__':
    main()
